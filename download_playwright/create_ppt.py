from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_strategy_slide():
    # 1. 创建演示文稿
    prs = Presentation()
    
    # 使用空白版式 (通常索引为6是空白页)
    slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(slide_layout)

    # --- 2. 添加标题 ---
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "年度经营目标与四步走战略路径"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102) # 深蓝色

    # --- 3. 添加左侧核心年度目标 ---
    goal_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(3), Inches(2))
    tf_goal = goal_box.text_frame
    
    p1 = tf_goal.add_paragraph()
    p1.text = "🏆 年度冲刺目标"
    p1.font.size = Pt(18)
    p1.font.color.rgb = RGBColor(0, 51, 102)
    
    p2 = tf_goal.add_paragraph()
    p2.text = "13 万+ 台"
    p2.font.size = Pt(54) # 超大字体
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255, 102, 0) # 橙色强调
    
    p3 = tf_goal.add_paragraph()
    p3.text = "Target Achievement"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(128, 128, 128)

    # --- 4. 绘制阶梯状的四个阶段 (使用形状) ---
    # 定义四个阶段的数据：(标题, 时间, 策略, 目标, 颜色)
    steps = [
        ("第一阶段：爆款驱动", "12月 - 2月", "自营裂变 · 养权重", "目标：1.5 万台", RGBColor(189, 215, 238)), # 浅蓝
        ("第二阶段：新品驱动", "3月 - 6月", "新品分销 · 618爆发", "目标：6 万台", RGBColor(155, 194, 230)), # 中蓝
        ("第三阶段：新商驱动", "7月 - 9月", "分销拓新 · 蓄水双11", "目标：8.5 万台", RGBColor(46, 117, 182)), # 深蓝
        ("第四阶段：大促驱动", "10月 - 12月", "全域共振 · 饱和攻击", "冲刺：13 万+ 台", RGBColor(31, 78, 121))  # 最深蓝
    ]

    # 起始位置
    start_x = 3.0
    start_y = 5.0
    width = 2.2
    height = 1.8
    step_rise = 1.0 # 每一级台阶上升的高度
    x_offset = 1.6  # 每一级台阶向右的位移

    for i, (title, time, strategy, target, color) in enumerate(steps):
        # 计算位置：形成阶梯状
        left = Inches(start_x + (i * x_offset))
        top = Inches(start_y - (i * step_rise))
        
        # 添加矩形背景
        shape = slide.shapes.add_shape(
            1, # MSO_SHAPE.RECTANGLE
            left, top, Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(255, 255, 255)

        # 在形状内添加文本
        tf = shape.text_frame
        tf.margin_top = Inches(0.1)
        tf.margin_left = Inches(0.1)
        
        # 阶段数字背景 (大背景数字，可选，这里简化为文字)
        p_num = tf.paragraphs[0]
        p_num.text = str(i + 1)
        p_num.font.size = Pt(40)
        p_num.font.bold = True
        p_num.font.color.rgb = RGBColor(255, 255, 255)
        p_num.alignment = PP_ALIGN.RIGHT # 数字放右上角装饰

        # 阶段内容
        p_title = tf.add_paragraph()
        p_title.text = title
        p_title.font.size = Pt(12)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(255, 255, 255)

        p_time = tf.add_paragraph()
        p_time.text = f"📅 {time}"
        p_time.font.size = Pt(10)
        p_time.font.color.rgb = RGBColor(230, 230, 230)

        p_strat = tf.add_paragraph()
        p_strat.text = f"⚙️ {strategy}"
        p_strat.font.size = Pt(10)
        p_strat.font.color.rgb = RGBColor(255, 255, 255)
        
        p_target = tf.add_paragraph()
        p_target.text = f"📈 {target}"
        p_target.font.size = Pt(11)
        p_target.font.bold = True
        p_target.font.color.rgb = RGBColor(255, 255, 0) # 黄色高亮目标

    # --- 5. 添加底部 Summary ---
    summary_box = slide.shapes.add_shape(
        1, # Rectangle
        Inches(0), Inches(6.8), Inches(10), Inches(0.7)
    )
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = RGBColor(0, 51, 102) # 深蓝底色
    
    tf_sum = summary_box.text_frame
    tf_sum.vertical_anchor = 3 # Middle
    p_sum = tf_sum.paragraphs[0]
    p_sum.text = "Summary: 通过“爆款-新品-新商-大促”四阶引擎轮动，实现从日销爬坡到大促爆发的全年稳定增长。"
    p_sum.font.size = Pt(14)
    p_sum.font.color.rgb = RGBColor(255, 255, 255)
    p_sum.alignment = PP_ALIGN.CENTER

    # 保存
    prs.save('年度全域增长战略.pptx')
    print("PPT生成成功：年度全域增长战略.pptx")

if __name__ == "__main__":
    create_strategy_slide()